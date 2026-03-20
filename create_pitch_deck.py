"""
Pitch deck generator for Claude Code Playground — AI Task Manager
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT_VIOLET= RGBColor(0x7C, 0x3A, 0xED)   # violet
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)   # cyan
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # emerald
ACCENT_MAGENTA=RGBColor(0xEC, 0x48, 0x99)   # magenta/pink
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xB0, 0xBA, 0xD4)
MID_GREY     = RGBColor(0x3A, 0x3A, 0x5C)
CARD_BG      = RGBColor(0x16, 0x16, 0x30)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    fill_slide_bg(slide, DARK_BG)
    return slide


def fill_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # rectangle = 1
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=20, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                          default_size=18, default_color=WHITE, default_bold=False,
                          align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """lines = list of (text, size, color, bold) tuples or plain strings"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, default_size, default_color, default_bold)
        text, size, color, bold = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txb


def accent_bar(slide, color, left=Inches(0), top=Inches(0),
               width=Inches(0.06), height=SLIDE_H):
    add_rect(slide, left, top, width, height, color)


def horiz_rule(slide, top, color=MID_GREY, left=Inches(0.6), width=None):
    if width is None:
        width = SLIDE_W - Inches(1.2)
    add_rect(slide, left, top, width, Inches(0.015), color)


def slide_number(slide, n, total=11):
    add_textbox(slide, f"{n:02d} / {total:02d}",
                SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.45),
                Inches(1.2), Inches(0.35),
                font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)


def tag_pill(slide, text, left, top, bg_color, text_color=WHITE):
    w, h = Inches(1.5), Inches(0.32)
    add_rect(slide, left, top, w, h, bg_color)
    add_textbox(slide, text, left, top, w, h,
                font_size=11, bold=True, color=text_color,
                align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_VIOLET)
accent_bar(s, ACCENT_CYAN, left=SLIDE_W - Inches(0.06))

# Background gradient rect (simulate)
add_rect(s, Inches(0.06), Inches(0), SLIDE_W - Inches(0.12), SLIDE_H,
         RGBColor(0x0D, 0x0D, 0x1A))

# Glow circle (decorative)
add_rect(s, Inches(8.5), Inches(-1), Inches(6), Inches(6),
         RGBColor(0x1A, 0x0D, 0x2E))

# Kicker
add_textbox(s, "AI HACKATHON  ·  DEMO PROJECT",
            Inches(0.9), Inches(1.3), Inches(8), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_CYAN,
            font_name="Consolas")

# Main title
add_textbox(s, "TASK // BOARD",
            Inches(0.9), Inches(2.0), Inches(9), Inches(1.5),
            font_size=64, bold=True, color=WHITE,
            font_name="Segoe UI Black")

# Subtitle
add_textbox(s, "An AI-native task manager where autonomous agents think,\ncreate, and complete work — in real time.",
            Inches(0.9), Inches(3.6), Inches(8.5), Inches(1.2),
            font_size=20, color=LIGHT_GREY)

# Tech pills
pill_data = [
    ("SAP CAP Java",  ACCENT_VIOLET, Inches(0.9)),
    ("LangGraph4j",   ACCENT_CYAN,   Inches(2.55)),
    ("React 19",      ACCENT_GREEN,  Inches(4.2)),
    ("Claude API",    ACCENT_MAGENTA,Inches(5.7)),
]
for label, color, left in pill_data:
    tag_pill(s, label, left, Inches(5.1), color)

slide_number(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_MAGENTA)

add_textbox(s, "THE PROBLEM", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_MAGENTA, font_name="Consolas")
add_textbox(s, "Task managers are passive tools.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.9),
            font_size=40, bold=True, color=WHITE)
horiz_rule(s, Inches(2.05))

problems = [
    ("🧍 Human-dependent",  "Every task requires a person to create it, track it, and close it."),
    ("⏳ Stale by default",  "Work boards go out of date the moment you stop updating them."),
    ("🔇 No real-time sync", "Multiple users and services see inconsistent views of the same data."),
]
for i, (title, desc) in enumerate(problems):
    top = Inches(2.4 + i * 1.4)
    add_rect(s, Inches(0.85), top, Inches(11.6), Inches(1.2), CARD_BG)
    add_textbox(s, title, Inches(1.1), top + Inches(0.1), Inches(4), Inches(0.45),
                font_size=18, bold=True, color=ACCENT_MAGENTA)
    add_textbox(s, desc, Inches(1.1), top + Inches(0.5), Inches(11), Inches(0.55),
                font_size=16, color=LIGHT_GREY)

slide_number(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_GREEN)

add_textbox(s, "THE SOLUTION", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_GREEN, font_name="Consolas")
add_textbox(s, "A task board that runs itself.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.9),
            font_size=40, bold=True, color=WHITE)
horiz_rule(s, Inches(2.05))

add_textbox(s,
    "TASK // BOARD combines an OData V4 REST backend with three autonomous AI agents "
    "that continuously generate tasks, complete them, and broadcast changes to every "
    "connected browser in real time — all without a single human click.",
    Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2),
    font_size=17, color=LIGHT_GREY)

features = [
    (ACCENT_VIOLET, "🤖 Self-operating",    "Agents create & complete tasks 24/7 without human input."),
    (ACCENT_CYAN,   "⚡ Real-time sync",    "Server-Sent Events push changes instantly to all clients."),
    (ACCENT_GREEN,  "💬 Chat interface",    "Natural language commands via a floating AI assistant."),
    (ACCENT_MAGENTA,"📊 Live analytics",   "Stats dashboard: avg duration, charts, full task history."),
]
for i, (color, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.85 + col * 6.1)
    top  = Inches(3.55 + row * 1.5)
    add_rect(s, left, top, Inches(5.8), Inches(1.3), CARD_BG)
    add_rect(s, left, top, Inches(0.06), Inches(1.3), color)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.1), Inches(5.4), Inches(0.45),
                font_size=17, bold=True, color=color)
    add_textbox(s, desc, left + Inches(0.2), top + Inches(0.55), Inches(5.4), Inches(0.6),
                font_size=14, color=LIGHT_GREY)

slide_number(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CYAN)

add_textbox(s, "ARCHITECTURE", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_CYAN, font_name="Consolas")
add_textbox(s, "Three layers. One living system.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.8),
            font_size=36, bold=True, color=WHITE)
horiz_rule(s, Inches(2.0))

# Three layer boxes
layers = [
    (ACCENT_VIOLET, "FRONTEND", "React 19 + Vite",
     "• Framer Motion animations\n• EventSource (SSE) listener\n• OData V4 fetch calls\n• Chat modal + Stats dashboard"),
    (ACCENT_CYAN,   "CAP BACKEND", "SAP CAP Java 2.6 + Spring Boot",
     "• OData V4 CRUD endpoints\n• suggestTask & chatTask actions\n• SSE controller + emitter registry\n• H2 (dev) / PostgreSQL (prod)"),
    (ACCENT_GREEN,  "AI AGENTS", "LangGraph4j + LangChain4j",
     "• TaskAgent — generates tasks every 10s\n• CompleterAgent — closes tasks (5–15s jitter)\n• MonitorAgent — detects changes every 2s\n• Claude API via Anthropic proxy"),
]
for i, (color, title, subtitle, body) in enumerate(layers):
    left = Inches(0.55 + i * 4.2)
    top  = Inches(2.3)
    w    = Inches(4.0)
    h    = Inches(4.5)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, w, Inches(0.06), color)
    add_textbox(s, title, left + Inches(0.15), top + Inches(0.18), w - Inches(0.3), Inches(0.4),
                font_size=16, bold=True, color=color, font_name="Consolas")
    add_textbox(s, subtitle, left + Inches(0.15), top + Inches(0.65), w - Inches(0.3), Inches(0.4),
                font_size=13, italic=True, color=LIGHT_GREY)
    horiz_rule(s, top + Inches(1.1), color=MID_GREY, left=left + Inches(0.15), width=w - Inches(0.3))
    add_textbox(s, body, left + Inches(0.15), top + Inches(1.25), w - Inches(0.3), Inches(3.0),
                font_size=13, color=LIGHT_GREY)

# Arrows between boxes
for arrow_left in [Inches(4.55), Inches(8.75)]:
    add_textbox(s, "⟷", arrow_left, Inches(4.3), Inches(0.4), Inches(0.5),
                font_size=22, color=MID_GREY, align=PP_ALIGN.CENTER)

slide_number(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE THREE AGENTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_VIOLET)

add_textbox(s, "THE AGENTS", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_VIOLET, font_name="Consolas")
add_textbox(s, "Autonomous. Parallel. Relentless.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.8),
            font_size=36, bold=True, color=WHITE)
horiz_rule(s, Inches(2.0))

agents = [
    (ACCENT_VIOLET, "🧠  TASK AGENT", "every 10 seconds",
     ["Brainstorms a new task with Claude",
      "Refines the title to start with 'AI '",
      "POSTs it to the OData endpoint",
      "Three-node LangGraph4j StateGraph"]),
    (ACCENT_CYAN,   "✅  COMPLETER AGENT", "every 5 – 15 seconds (jitter)",
     ["Fetches all pending tasks",
      "Picks one at random",
      "Marks it complete + sets completedAt",
      "Random delay prevents thundering herd"]),
    (ACCENT_GREEN,  "👁  MONITOR AGENT", "every 2 seconds",
     ["Polls the task list via OData GET",
      "Diffs the response against last snapshot",
      "Broadcasts 'tasks_changed' SSE event",
      "Triggers instant UI refresh in all browsers"]),
]
for i, (color, title, freq, bullets) in enumerate(agents):
    left = Inches(0.55 + i * 4.2)
    top  = Inches(2.35)
    w, h = Inches(4.0), Inches(4.6)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, w, Inches(0.06), color)
    add_textbox(s, title, left + Inches(0.15), top + Inches(0.18), w - Inches(0.3), Inches(0.45),
                font_size=16, bold=True, color=color)
    add_textbox(s, freq, left + Inches(0.15), top + Inches(0.7), w - Inches(0.3), Inches(0.35),
                font_size=12, italic=True, color=ACCENT_CYAN, font_name="Consolas")
    horiz_rule(s, top + Inches(1.1), color=MID_GREY, left=left + Inches(0.15), width=w - Inches(0.3))
    for j, bullet in enumerate(bullets):
        add_textbox(s, f"→  {bullet}", left + Inches(0.15), top + Inches(1.25 + j * 0.72),
                    w - Inches(0.3), Inches(0.65), font_size=13, color=LIGHT_GREY)

slide_number(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — REAL-TIME SYNC FLOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CYAN)

add_textbox(s, "REAL-TIME FLOW", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_CYAN, font_name="Consolas")
add_textbox(s, "From agent action to browser refresh in < 2 seconds.",
            Inches(0.9), Inches(1.1), Inches(11), Inches(0.8),
            font_size=34, bold=True, color=WHITE)
horiz_rule(s, Inches(2.0))

steps = [
    (ACCENT_VIOLET, "1", "Agent writes",     "CompleterAgent PATCHes a task\nvia OData V4 endpoint"),
    (ACCENT_CYAN,   "2", "DB changes",       "CAP persists the update\nin H2 / PostgreSQL"),
    (ACCENT_GREEN,  "3", "Monitor detects",  "MonitorAgent diffs the\nnew response vs snapshot"),
    (ACCENT_MAGENTA,"4", "SSE broadcast",    "SseEmitterRegistry sends\n'tasks_changed' to all clients"),
    (WHITE,         "5", "UI refreshes",     "EventSource listener calls\nloadTasks() — no polling"),
]
for i, (color, num, title, desc) in enumerate(steps):
    left = Inches(0.55 + i * 2.48)
    top  = Inches(2.5)
    w, h = Inches(2.3), Inches(3.8)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, w, Inches(0.06), color)
    add_textbox(s, num, left + Inches(0.1), top + Inches(0.18), Inches(0.6), Inches(0.5),
                font_size=30, bold=True, color=color, font_name="Segoe UI Black")
    add_textbox(s, title, left + Inches(0.1), top + Inches(0.85), w - Inches(0.2), Inches(0.4),
                font_size=15, bold=True, color=WHITE)
    add_textbox(s, desc, left + Inches(0.1), top + Inches(1.35), w - Inches(0.2), Inches(1.8),
                font_size=13, color=LIGHT_GREY)
    if i < 4:
        add_textbox(s, "→", left + w, top + Inches(1.7), Inches(0.18), Inches(0.4),
                    font_size=20, color=MID_GREY, align=PP_ALIGN.CENTER)

slide_number(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_MAGENTA)

add_textbox(s, "AI INTEGRATION", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_MAGENTA, font_name="Consolas")
add_textbox(s, "Claude does the thinking. LangChain4j does the wiring.",
            Inches(0.9), Inches(1.1), Inches(11.5), Inches(0.8),
            font_size=33, bold=True, color=WHITE)
horiz_rule(s, Inches(2.0))

integrations = [
    (ACCENT_VIOLET, "SUGGEST TASK",
     "User provides a prompt hint → Claude returns a creative\ntask title + description via a typed LangChain4j interface.",
     "POST /odata/v4/api/suggestTask"),
    (ACCENT_CYAN, "CHAT ASSISTANT",
     "User sends a natural language message + the full task list as\ncontext → Claude returns a structured action (create / complete\n/ reopen / delete / deleteAll) + parameters.",
     "POST /odata/v4/api/chatTask"),
    (ACCENT_GREEN, "AUTONOMOUS GENERATION",
     "TaskAgent uses a two-step LangGraph4j StateGraph: first\nbrainstorm a task, then refine the title to start with 'AI '.",
     "LangGraph4j StateGraph (3 nodes)"),
    (ACCENT_MAGENTA, "STRUCTURED OUTPUT",
     "All three use LangChain4j typed interfaces (@SystemMessage /\n@UserMessage) — no manual JSON parsing, just Java records.",
     "LangChain4j 1.0.0-beta3"),
]
for i, (color, title, desc, tag) in enumerate(integrations):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.35)
    top  = Inches(2.4 + row * 2.3)
    w, h = Inches(6.1), Inches(2.1)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, Inches(0.06), h, color)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.12), w - Inches(0.3), Inches(0.4),
                font_size=15, bold=True, color=color, font_name="Consolas")
    add_textbox(s, desc, left + Inches(0.2), top + Inches(0.6), w - Inches(0.3), Inches(0.95),
                font_size=13, color=LIGHT_GREY)
    add_textbox(s, tag, left + Inches(0.2), top + Inches(1.65), w - Inches(0.3), Inches(0.35),
                font_size=11, color=MID_GREY, italic=True, font_name="Consolas")

slide_number(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_GREEN)

add_textbox(s, "TECH STACK", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_GREEN, font_name="Consolas")
add_textbox(s, "Modern. Enterprise-ready. Fully open.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.7),
            font_size=36, bold=True, color=WHITE)
horiz_rule(s, Inches(1.9))

categories = [
    (ACCENT_VIOLET, "BACKEND",
     ["SAP CAP Java 2.6.0", "Spring Boot 3.2.5", "OData V4 adapter", "H2 (dev) / PostgreSQL (prod)"]),
    (ACCENT_CYAN, "AI / AGENTS",
     ["LangChain4j 1.0.0-beta3", "LangGraph4j 1.8.10", "Anthropic Claude API", "claude-haiku-4-5 (default)"]),
    (ACCENT_GREEN, "FRONTEND",
     ["React 19.2.4", "Vite 5.4.21", "Framer Motion 12", "Lucide React icons"]),
    (ACCENT_MAGENTA, "DEPLOYMENT",
     ["Docker multi-stage build", "Railway (backend)", "Vercel (frontend)", "Java 21 JRE runtime"]),
]
for i, (color, cat, items) in enumerate(categories):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.35)
    top  = Inches(2.25 + row * 2.3)
    w, h = Inches(6.1), Inches(2.1)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, w, Inches(0.06), color)
    add_textbox(s, cat, left + Inches(0.2), top + Inches(0.18), w, Inches(0.4),
                font_size=14, bold=True, color=color, font_name="Consolas")
    for j, item in enumerate(items):
        add_textbox(s, f"▸  {item}", left + Inches(0.2), top + Inches(0.7 + j * 0.34),
                    w - Inches(0.3), Inches(0.32), font_size=14, color=LIGHT_GREY)

slide_number(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KEY DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_VIOLET)

add_textbox(s, "DESIGN DECISIONS", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_VIOLET, font_name="Consolas")
add_textbox(s, "Engineering choices that matter.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.7),
            font_size=36, bold=True, color=WHITE)
horiz_rule(s, Inches(1.9))

decisions = [
    (ACCENT_CYAN,    "SSE over WebSockets",
     "Simpler, HTTP/1.1-compatible, unidirectional push — perfect for read-heavy broadcast without bi-directional overhead."),
    (ACCENT_VIOLET,  "LangGraph4j for agents",
     "Composable StateGraph nodes enable multi-step reasoning (generate → refine → post) with clean separation of concerns."),
    (ACCENT_GREEN,   "OData V4 for CRUD",
     "CAP auto-generates filter, select, and expand operations. Standard protocol enables future SAP BTP integration."),
    (ACCENT_MAGENTA, "Random jitter in CompleterAgent",
     "5–15s random delay prevents thundering-herd bursts when multiple instances run in parallel on the same schedule."),
    (ACCENT_CYAN,    "Typed LangChain4j interfaces",
     "Java records replace manual JSON parsing. Compiler enforces response shape. Simpler, safer, more maintainable."),
    (ACCENT_VIOLET,  "HTTP/1.1 for Claude API calls",
     "SAP Hyperspace proxy rejects HTTP/2. Explicit HttpClient.Version.HTTP_1_1 avoids silent protocol negotiation failure."),
]
for i, (color, title, desc) in enumerate(decisions):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.35)
    top  = Inches(2.2 + row * 1.65)
    w, h = Inches(6.1), Inches(1.5)
    add_rect(s, left, top, w, h, CARD_BG)
    add_rect(s, left, top, Inches(0.06), h, color)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.1), w - Inches(0.3), Inches(0.4),
                font_size=15, bold=True, color=color)
    add_textbox(s, desc, left + Inches(0.2), top + Inches(0.55), w - Inches(0.3), Inches(0.85),
                font_size=12, color=LIGHT_GREY)

slide_number(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CYAN)

add_textbox(s, "DEPLOYMENT", Inches(0.9), Inches(0.55), Inches(10), Inches(0.5),
            font_size=13, bold=True, color=ACCENT_CYAN, font_name="Consolas")
add_textbox(s, "Production-ready on day one.", Inches(0.9), Inches(1.1), Inches(11), Inches(0.7),
            font_size=36, bold=True, color=WHITE)
horiz_rule(s, Inches(1.9))

# Railway box
add_rect(s, Inches(0.6), Inches(2.3), Inches(5.9), Inches(4.5), CARD_BG)
add_rect(s, Inches(0.6), Inches(2.3), Inches(5.9), Inches(0.06), ACCENT_VIOLET)
add_textbox(s, "🚂  RAILWAY — Backend", Inches(0.8), Inches(2.48), Inches(5.5), Inches(0.45),
            font_size=17, bold=True, color=ACCENT_VIOLET)
be_items = [
    "Docker multi-stage build (Node 22 + Maven → JRE 21)",
    "Starts: java -Dspring.profiles.active=prod -jar app.jar",
    "Health check: GET /odata/v4/api/",
    "Restart policy: ON_FAILURE",
    "Env vars: ANTHROPIC_API_KEY, ANTHROPIC_PROXY_URL, ANTHROPIC_MODEL",
]
for j, item in enumerate(be_items):
    add_textbox(s, f"▸  {item}", Inches(0.8), Inches(3.1 + j * 0.55), Inches(5.5), Inches(0.5),
                font_size=13, color=LIGHT_GREY)

# Vercel box
add_rect(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(4.5), CARD_BG)
add_rect(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(0.06), ACCENT_CYAN)
add_textbox(s, "▲  VERCEL — Frontend", Inches(7.05), Inches(2.48), Inches(5.5), Inches(0.45),
            font_size=17, bold=True, color=ACCENT_CYAN)
fe_items = [
    "npm install && npm run build (Vite → dist/)",
    "Static hosting from dist/",
    "VITE_API_URL env var points to Railway backend",
    "Zero-config framework detection for Vite",
    "Global CDN, automatic HTTPS",
]
for j, item in enumerate(fe_items):
    add_textbox(s, f"▸  {item}", Inches(7.05), Inches(3.1 + j * 0.55), Inches(5.5), Inches(0.5),
                font_size=13, color=LIGHT_GREY)

slide_number(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING / THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_VIOLET)
accent_bar(s, ACCENT_CYAN, left=SLIDE_W - Inches(0.06))

add_rect(s, Inches(0.06), Inches(0), SLIDE_W - Inches(0.12), SLIDE_H,
         RGBColor(0x0D, 0x0D, 0x1A))

add_textbox(s, "THANK YOU",
            Inches(0.9), Inches(1.6), Inches(11), Inches(1.4),
            font_size=64, bold=True, color=WHITE, font_name="Segoe UI Black")

add_textbox(s, "TASK // BOARD — AI-native task management, built in a day.",
            Inches(0.9), Inches(3.1), Inches(11), Inches(0.7),
            font_size=22, color=LIGHT_GREY)

# Bottom tag strip
pills = [
    ("SAP CAP Java",   ACCENT_VIOLET, Inches(0.9)),
    ("LangGraph4j",    ACCENT_CYAN,   Inches(2.55)),
    ("React 19",       ACCENT_GREEN,  Inches(4.2)),
    ("Claude API",     ACCENT_MAGENTA,Inches(5.7)),
    ("Railway",        ACCENT_VIOLET, Inches(7.2)),
    ("Vercel",         ACCENT_CYAN,   Inches(8.85)),
]
for label, color, left in pills:
    tag_pill(s, label, left, Inches(4.2), color)

add_textbox(s, "github.tools.sap/CALMBuild/claude-code-playground",
            Inches(0.9), Inches(5.3), Inches(11), Inches(0.4),
            font_size=13, color=MID_GREY, italic=True, font_name="Consolas")

slide_number(s, 11)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/Users/d047480/git/claude-code-playground/TaskBoard_PitchDeck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
